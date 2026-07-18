"""Render a GrowthReview as a slide deck (python-pptx).

Slides carry the same Figure-rendered strings as the markdown memo — one
source of truth, two output formats.
"""

from __future__ import annotations

from pathlib import Path
from typing import Any

import pandas as pd
from pptx import Presentation
from pptx.util import Inches, Pt

from growth_lab.reporting.growth_review import GrowthReview

_KPI_SLUGS = ("spend", "signups", "paid_signups", "revenue", "cac", "fraud_rate")


def _add_text_slide(prs: Any, title: str, lines: list[str]) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[5])  # title-only layout
    slide.shapes.title.text = title
    box = slide.shapes.add_textbox(Inches(0.5), Inches(1.5), Inches(9.0), Inches(5.5))
    frame = box.text_frame
    frame.word_wrap = True
    for i, line in enumerate(lines):
        paragraph = frame.paragraphs[0] if i == 0 else frame.add_paragraph()
        paragraph.text = line
        paragraph.font.size = Pt(14)


def _frame_lines(frame: pd.DataFrame) -> list[str]:
    def cell(value: object) -> str:
        if isinstance(value, float):
            return f"{value:,.2f}"
        if isinstance(value, int):
            return f"{value:,}"
        return str(value)

    header = "  |  ".join(str(c) for c in frame.columns)
    rows = ["  |  ".join(cell(v) for v in row) for row in frame.itertuples(index=False, name=None)]
    return [header, *rows]


def export_pptx(review: GrowthReview, path: Path) -> Path:
    prs = Presentation()
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = "Meridian growth review"
    title_slide.placeholders[1].text = f"Data through {review.as_of}"

    f = review.figures.get
    kpi_lines = [
        f"{slug.replace('_', ' ').title()}: {f('kpi_' + slug).render()} "
        f"({f('kpi_' + slug + '_wow').render()} WoW)"
        for slug in _KPI_SLUGS
    ]
    _add_text_slide(prs, "KPIs — trailing week", kpi_lines)

    channel_frame = f("channel_week").value
    assert isinstance(channel_frame, pd.DataFrame)
    _add_text_slide(prs, "Channel performance", _frame_lines(channel_frame))

    plans = f("ltv_plans").value
    assert isinstance(plans, pd.DataFrame)
    _add_text_slide(prs, "Customer lifetime value", _frame_lines(plans))

    _add_text_slide(
        prs,
        "Outlook and monitoring",
        [
            f"Expected revenue, next {f('forecast_horizon').render()} days: "
            f"{f('forecast_total').render()}",
            f"Anomalous revenue days flagged: {f('anomaly_count').render()}",
        ],
    )

    path.parent.mkdir(parents=True, exist_ok=True)
    prs.save(str(path))
    return path
