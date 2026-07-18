"""Phase 5 gate: every number in the memo traces to a sourced figure.

The extraction is exhaustive: every numeric token in the rendered markdown
must appear in some figure's rendered output (dates and the ISO date tokens
excepted). If someone hand-types a number into the template, this fails.
"""

from __future__ import annotations

import ast
import re
from pathlib import Path

import duckdb
import pytest
from pptx import Presentation

from growth_lab.reporting import (
    Figure,
    FigureRegistry,
    GrowthReview,
    build_weekly_review,
    export_pptx,
)

DATE_TOKEN = re.compile(r"\b\d{4}-\d{2}-\d{2}\b")
NUMBER_TOKEN = re.compile(r"-?\d[\d,]*\.?\d*")


@pytest.fixture(scope="module")
def review(warehouse: duckdb.DuckDBPyConnection) -> GrowthReview:
    return build_weekly_review(warehouse)


def test_every_number_in_memo_is_figure_backed(review: GrowthReview) -> None:
    memo = DATE_TOKEN.sub("", review.markdown)
    memo_numbers = set(NUMBER_TOKEN.findall(memo))

    backed: set[str] = set()
    for figure in review.figures.all():
        rendered = DATE_TOKEN.sub("", figure.render())
        backed.update(NUMBER_TOKEN.findall(rendered))
        backed.update(NUMBER_TOKEN.findall(DATE_TOKEN.sub("", figure.source)))

    orphans = memo_numbers - backed
    assert not orphans, f"numbers with no backing figure: {sorted(orphans)}"


def test_memo_has_all_sections(review: GrowthReview) -> None:
    for heading in (
        "# Meridian growth review",
        "## KPIs, trailing week",
        "## Channel performance",
        "## Customer lifetime value",
        "## Outlook and monitoring",
        "## Provenance",
    ):
        assert heading in review.markdown, heading


def test_every_figure_has_provenance(review: GrowthReview) -> None:
    for figure in review.figures.all():
        assert figure.source.strip(), figure.slug
    # KPI figures must carry the actual semantic-layer SQL
    assert "SELECT" in review.figures.get("kpi_cac").source
    assert "mart_daily_channel" in review.figures.get("kpi_cac").source


def test_review_is_deterministic(warehouse: duckdb.DuckDBPyConnection) -> None:
    first = build_weekly_review(warehouse)
    second = build_weekly_review(warehouse)
    assert first.markdown == second.markdown


def test_registry_rejects_duplicates_and_missing_sources() -> None:
    registry = FigureRegistry()
    registry.add(Figure(slug="a", title="a", value=1.0, source="somewhere"))
    with pytest.raises(ValueError, match="duplicate"):
        registry.add(Figure(slug="a", title="again", value=2.0, source="somewhere"))
    with pytest.raises(ValueError, match="provenance"):
        Figure(slug="b", title="b", value=1.0, source="   ")
    with pytest.raises(KeyError):
        registry.get("nope")


def test_pptx_export_round_trips(review: GrowthReview, tmp_path: Path) -> None:
    path = export_pptx(review, tmp_path / "weekly_review.pptx")
    deck = Presentation(str(path))
    assert len(deck.slides) == 5
    all_text = " ".join(
        shape.text_frame.text
        for slide in deck.slides
        for shape in slide.shapes
        if shape.has_text_frame
    )
    assert "Meridian growth review" in all_text
    assert review.figures.get("kpi_revenue").render() in all_text


def test_dashboard_app_parses() -> None:
    """Streamlit is an optional extra; CI verifies the app is at least
    syntactically sound without importing it."""
    source = (Path(__file__).resolve().parents[1] / "dashboard" / "app.py").read_text()
    ast.parse(source)
    assert "st.tabs" in source
